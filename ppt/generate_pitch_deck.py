from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
SVG_LOGO_PATH = HERE / "CLAi-logo.svg"
LOGO_PNG_PATH = HERE / "CLAi-logo.png"
OUTPUT_PPTX_PATH = HERE / "CodeLaunchAI-Pitch-Deck.pptx"


def ensure_logo_png() -> None:
    if not SVG_LOGO_PATH.exists():
        raise FileNotFoundError(f"Missing logo SVG: {SVG_LOGO_PATH}")

    if LOGO_PNG_PATH.exists():
        return

    # Prefer CairoSVG if it's usable, but fall back to a pure-Python renderer on Windows
    # environments missing native cairo DLLs.
    try:
        import cairosvg  # type: ignore

        cairosvg.svg2png(
            url=str(SVG_LOGO_PATH),
            write_to=str(LOGO_PNG_PATH),
            output_width=320,
            output_height=90,
        )
        return
    except Exception:
        pass

    try:
        from PIL import Image, ImageColor, ImageDraw, ImageFont
    except Exception as exc:  # pragma: no cover
        raise RuntimeError(
            "Logo conversion failed (CairoSVG unavailable) and Pillow is not installed. "
            "Install Pillow: `pip install pillow`."
        ) from exc

    width, height = 320, 90
    img = Image.new("RGBA", (width, height), (0, 0, 0, 0))

    def lerp(a: int, b: int, t: float) -> int:
        return int(round(a + (b - a) * t))

    def make_linear_gradient(size, start_hex: str, end_hex: str):
        w, h = size
        start = ImageColor.getrgb(start_hex)
        end = ImageColor.getrgb(end_hex)
        grad = Image.new("RGBA", (w, h))
        px = grad.load()
        for x in range(w):
            t = x / (w - 1) if w > 1 else 0.0
            r = lerp(start[0], end[0], t)
            g = lerp(start[1], end[1], t)
            b = lerp(start[2], end[2], t)
            for y in range(h):
                px[x, y] = (r, g, b, 255)
        return grad

    # Icon
    icon_grad = make_linear_gradient((72, 72), "#90F7EC", "#32B2FF")

    # Soft circle background (12% opacity)
    circle_mask = Image.new("L", (72, 72), 0)
    ImageDraw.Draw(circle_mask).ellipse((0, 0, 71, 71), fill=int(255 * 0.12))
    img.paste(icon_grad, (8, 9), circle_mask)

    # Diamond
    diamond_mask = Image.new("L", (72, 72), 0)
    ImageDraw.Draw(diamond_mask).polygon([(36, 11), (46, 60), (36, 52), (26, 60)], fill=255)
    img.paste(icon_grad, (8, 9), diamond_mask)

    # Inner circles
    draw = ImageDraw.Draw(img)
    cx, cy = 42, 45
    draw.ellipse((cx - 7, cy - 7, cx + 7, cy + 7), fill=(255, 255, 255, int(255 * 0.85)))
    small_grad = make_linear_gradient((8, 8), "#90F7EC", "#32B2FF")
    small_mask = Image.new("L", (8, 8), 0)
    ImageDraw.Draw(small_mask).ellipse((0, 0, 7, 7), fill=255)
    img.paste(small_grad, (cx - 4, cy - 4), small_mask)

    # Text "CLAi" with gradient fill
    text_grad = make_linear_gradient((width, height), "#04C7C1", "#286BFF")
    text_mask = Image.new("L", (width, height), 0)
    mask_draw = ImageDraw.Draw(text_mask)

    font = None
    for candidate in ("Inter.ttf", "arial.ttf", "segoeui.ttf"):
        try:
            font = ImageFont.truetype(candidate, 54)
            break
        except Exception:
            continue
    if font is None:
        font = ImageFont.load_default()

    # Approximate placement (matches the SVG roughly)
    mask_draw.text((80, 18), "CLAi", font=font, fill=255)
    img = Image.composite(text_grad, img, text_mask)

    img.save(LOGO_PNG_PATH)


slides = [
    {
        "title": "AI-Powered Web App Builder",
        "subtitle": "From Idea to Paid SaaS – Build Real Apps for $10/Month",
        "footer": "Saravanan Rajagopalan / CodeLaunchAI\n2025-12-29",
        "logo": True,
        "center": True,
    },
    {
        "title": "Problem Statement",
        "bullets": [
            "Learning web development is slow and fragmented.",
            "Existing no-code / AI tools are expensive or limited.",
            "Students and indie developers struggle to build production-ready apps with payments.",
        ],
        "footer": "Complexity • Time Lost • Money Wasted",
    },
    {
        "title": "Solution Overview",
        "bullets": [
            "AI generates full-stack apps (frontend + backend)",
            "Built-in Stripe payments & subscriptions",
            "Credit-based AI usage system (cost-efficient)",
            "Editable, deployable, learning-focused code",
        ],
    },
    {
        "title": "Product Demo / Screenshots",
        "bullets": [
            "Landing page screenshot (student-friendly UI)",
            "AI app generation workflow",
            "Stripe integration example",
            "Code export / deployment interface",
        ],
        "note": "Insert relevant screenshots or mockups here.",
    },
    {
        "title": "Target Market (TAM / SAM / SOM)",
        "bullets": [
            "TAM: 70M global students, bootcamp attendees, self-taught devs",
            "SAM: 700,000 reachable users in 1–3 years",
            "SOM: 7,000 paying users early phase",
        ],
        "note": "Show funnel/bar chart comparison.",
    },
    {
        "title": "Competitor Comparison Matrix",
        "table": {
            "columns": [
                "Feature",
                "Our Product",
                "Replit",
                "Lovable AI",
                "Bubble",
                "GitHub Copilot",
                "UI-Only AI Tools",
            ],
            "rows": [
                ["Full-stack", "✅", "Partial", "Partial", "Yes", "No", "No"],
                ["Stripe subs", "Built-in", "Manual", "No", "Complex", "No", "No"],
                ["Editable code", "✅", "Yes", "Limited", "No", "Yes", "Limited"],
                ["Student-friendly", "$10", "$10–$20", "$20+", "$29+", "$10", "$20+"],
                ["AI usage control", "Credit", "Unlimited", "Limited", "N/A", "Unlimited", "Limited"],
            ],
        },
    },
    {
        "title": "Revenue Model / Pricing",
        "bullets": [
            "Student Plan: $10/month (200 AI credits)",
            "Pro Plan: $25/month (1,000 AI credits)",
            "Upsell: Add-on credits, templates, enterprise licensing",
            "SaaS recurring revenue focus",
        ],
    },
    {
        "title": "Unit Economics",
        "bullets": [
            "CAC: $5–$10 per user",
            "Monthly revenue: $10 (student plan)",
            "Gross margin: ~70%",
            "Payback period: 1–2 months",
            "LTV: $70–$100 (7–10 months retention)",
        ],
        "note": "Insert bar chart or KPIs visualization.",
    },
    {
        "title": "Marketing & Customer Acquisition",
        "bullets": [
            "Channels: YouTube tutorials, Reddit communities, Twitter/X, Google Ads",
            "Low-cost, student & dev-focused campaigns",
            "Organic growth via tutorials and templates",
        ],
        "note": "Funnel visualization (awareness → signup → paid subscriber)",
    },
    {
        "title": "Roadmap / Future Features",
        "bullets": [
            "Template marketplace",
            "Team collaboration / multi-user support",
            "GitHub sync & version control",
            "Mobile support / PWA",
            "AI model improvements (code refactoring, optimization)",
        ],
    },
    {
        "title": "Risks & Mitigation",
        "bullets": [
            "High AI costs → Credit limits & efficient prompt handling",
            "Student churn → Learning paths + engagement",
            "Payment failures/fraud → Stripe integration & monitoring",
            "Abuse → Rate-limiting + AI usage logs",
        ],
    },
    {
        "title": "Build Real SaaS Apps from Idea to Deployment",
        "subtitle": "Contact: CodeLaunchAI.ai",
        "footer": "Let’s talk – Investors & early partners welcome!",
        "center": True,
        "logo": True,
    },
]


def add_logo(slide, left, top, width=Inches(2.0)):
    slide.shapes.add_picture(str(LOGO_PNG_PATH), left, top, width=width)


def add_watermark(slide):
    left = Inches(8)
    top = Inches(6.8)
    width = Inches(1.1)
    slide.shapes.add_picture(str(LOGO_PNG_PATH), left, top, width=width)


def set_textbox(textbox, text: str, *, size: Pt, bold: bool = False, italic: bool = False, align=None, color=None):
    tf = textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.italic = italic
    if align is not None:
        p.alignment = align
    if color is not None:
        p.font.color.rgb = color


def main() -> int:
    ensure_logo_png()

    prs = Presentation()
    prs.slide_height = Inches(7.5)
    prs.slide_width = Inches(13.33)

    for idx, content in enumerate(slides):
        layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(layout)

        title_shape = slide.shapes.title
        title_shape.text = content["title"]
        tp = title_shape.text_frame.paragraphs[0]
        tp.font.bold = True
        tp.font.size = Pt(40)

        if content.get("subtitle"):
            subtitle = slide.shapes.add_textbox(Inches(2), Inches(1.8), Inches(10), Inches(1))
            set_textbox(subtitle, content["subtitle"], size=Pt(28), align=PP_ALIGN.CENTER)

        if content.get("bullets"):
            bulletbox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(4))
            tf = bulletbox.text_frame
            tf.clear()
            for i, item in enumerate(content["bullets"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = item
                p.font.size = Pt(24)
                p.level = 0

        if content.get("note"):
            notebox = slide.shapes.add_textbox(Inches(1.5), Inches(6.2), Inches(10), Inches(1))
            set_textbox(
                notebox,
                "Note: " + content["note"],
                size=Pt(15),
                italic=True,
                color=RGBColor(100, 100, 100),
            )

        if content.get("footer"):
            footerbox = slide.shapes.add_textbox(Inches(1), Inches(6.9), Inches(11.5), Inches(1))
            set_textbox(
                footerbox,
                content["footer"],
                size=Pt(15),
                color=RGBColor(70, 90, 140),
            )

        if content.get("table"):
            table_data = content["table"]
            rows, cols = len(table_data["rows"]) + 1, len(table_data["columns"])
            left, top, width, height = Inches(1.2), Inches(2.2), Inches(11), Inches(2.6)
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            for cidx, val in enumerate(table_data["columns"]):
                cell = table.cell(0, cidx)
                cell.text = val
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(16)

            for ridx, row in enumerate(table_data["rows"]):
                for cidx, val in enumerate(row):
                    cell = table.cell(ridx + 1, cidx)
                    cell.text = val
                    cell.text_frame.paragraphs[0].font.size = Pt(14)

        if content.get("center"):
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        if content.get("logo"):
            add_logo(slide, Inches(9.5), Inches(0.4), Inches(3.0))

        if idx != 0 and idx != len(slides) - 1:
            add_watermark(slide)

    prs.save(str(OUTPUT_PPTX_PATH))
    print(f"Pitch deck saved as {OUTPUT_PPTX_PATH}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
